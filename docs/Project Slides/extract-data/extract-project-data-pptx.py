"""
Extract project data directly from PowerPoint (.pptx) files
This is much more accurate than OCR since we get the actual text content.

Requirements:
- pip install python-pptx

Run with: python extract-project-data-pptx.py

Note: Script is in docs/Project Slides/extract-data/, paths are relative to project root
"""

import json
import re
import sys
from pathlib import Path
from typing import Dict, List, Optional
from datetime import datetime

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: Missing required library. Install with: pip install python-pptx")
    sys.exit(1)

# Base directory (project root)
BASE_DIR = Path(__file__).parent.parent.parent.parent
PROJECT_SLIDES_DIR = BASE_DIR / "docs" / "Project Slides"
PROJECTS_JSON = BASE_DIR / "data" / "projects.json"
LOG_FILE = BASE_DIR / "docs" / "Project Slides" / "extract-data" / "pptx_extraction_log.txt"
REPORT_FILE = BASE_DIR / "docs" / "Project Slides" / "pptx_extraction_report.md"

# Ensure log directory exists
LOG_FILE.parent.mkdir(parents=True, exist_ok=True)

def log(message: str):
    """Log message to file and console"""
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    log_message = f"[{timestamp}] {message}"
    print(log_message)
    with open(LOG_FILE, 'a', encoding='utf-8') as f:
        f.write(log_message + '\n')

def normalize_name(name: str) -> str:
    """Normalize project name for matching"""
    return re.sub(r'[^\w\s]', '', name.lower().strip())

def extract_text_from_pptx(pptx_path: Path) -> str:
    """Extract all text from PowerPoint presentation"""
    try:
        log(f"Processing PowerPoint: {pptx_path.name}")
        prs = Presentation(pptx_path)
        
        all_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            if slide_text:
                all_text.append(f"=== Slide {slide_num} ===")
                all_text.extend(slide_text)
                all_text.append("")  # Empty line between slides
        
        full_text = '\n'.join(all_text)
        log(f"Extracted {len(full_text)} characters of text from {len(prs.slides)} slides")
        return full_text
    except Exception as e:
        log(f"ERROR extracting text from {pptx_path}: {e}")
        return ''

def find_projects_in_category(text: str, category: str, existing_projects: List[Dict]) -> List[Dict]:
    """Find projects mentioned in text - looks for lines starting with 'Project'"""
    projects = []
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    
    # Debug: Log first few lines
    log(f"First 20 lines of text:")
    for i, line in enumerate(lines[:20]):
        log(f"  Line {i}: {line}")
    
    # Find all lines containing "Project"
    project_lines = []
    for i, line in enumerate(lines):
        if 'project' in line.lower():
            project_lines.append((i, line))
            log(f"  Found 'Project' line {i}: {line}")
    
    log(f"Found {len(project_lines)} 'Project' lines in text")
    
    # Try direct matching of project titles
    log("Trying direct project title matching...")
    text_lower = text.lower()
    text_no_spaces = text_lower.replace(' ', '').replace('&', '').replace(',', '').replace('-', '').replace(':', '')
    
    for project in existing_projects:
        project_title = project.get('title', '')
        if not project_title or any(p['title'] == project_title for p in projects):
            continue
        
        # Extract key words from project title
        title_words = project_title.split()
        key_words = [w for w in title_words if len(w) > 3 and w.lower() not in ['the', 'and', 'of', 'for', 'in', 'on', 'at', 'to', 'a', 'an']]
        
        if not key_words:
            key_words = [w for w in title_words if len(w) > 2]
        
        if not key_words:
            continue
        
        # Check if title (without spaces) appears in text
        title_no_spaces = ''.join(title_words).lower().replace('&', '').replace(',', '').replace('-', '').replace(':', '')
        
        if title_no_spaces and len(title_no_spaces) > 5 and title_no_spaces in text_no_spaces:
            # Find which line contains it
            for i, line in enumerate(lines):
                line_no_spaces = normalize_name(line).replace(' ', '').replace('&', '').replace(',', '').replace('-', '')
                if title_no_spaces in line_no_spaces:
                    projects.append({
                        'title': project_title,
                        'project': project,
                        'line': line,
                        'line_index': i,
                        'found_name': project_title,
                        'match_method': 'no_spaces_match'
                    })
                    log(f"  Direct no-spaces match: '{project_title}' in line {i}")
                    break
            continue
        
        # Check if enough key words appear
        matches = 0
        matched_words = []
        for word in key_words:
            word_lower = word.lower()
            word_no_spaces = word_lower.replace('-', '')
            if word_lower in text_lower or word_no_spaces in text_no_spaces:
                matches += 1
                matched_words.append(word)
        
        required_matches = min(2, len(key_words)) if len(key_words) > 1 else 1
        if matches >= required_matches:
            # Find which line contains the matches (prefer lines with "project")
            best_line_idx = None
            best_line = None
            for i, line in enumerate(lines):
                line_lower = line.lower()
                line_no_spaces = line_lower.replace(' ', '').replace('&', '').replace(',', '').replace('-', '')
                line_matches = sum(1 for word in matched_words if word.lower() in line_lower or word.lower().replace('-', '') in line_no_spaces)
                
                has_project = 'project' in line_lower
                if line_matches >= required_matches:
                    if best_line_idx is None or (has_project and 'project' not in (lines[best_line_idx].lower() if best_line_idx is not None else '')):
                        best_line_idx = i
                        best_line = line
                        if has_project:
                            break
            
            if best_line_idx is not None:
                projects.append({
                    'title': project_title,
                    'project': project,
                    'line': best_line,
                    'line_index': best_line_idx,
                    'found_name': project_title,
                    'match_method': 'keyword_match',
                    'matched_words': matched_words
                })
                log(f"  Direct keyword match: '{project_title}' (matched {matches}/{len(key_words)} words: {matched_words}) in line {best_line_idx}")
    
    log(f"Direct matching found {len(projects)} projects")
    
    # Also try matching extracted project names from "Project" lines
    for line_idx, line in project_lines:
        # Extract project name after "Project"
        project_parts = re.split(r'project\s*:?\s*', line, flags=re.IGNORECASE)
        for part in project_parts:
            if not part.strip():
                continue
            
            project_name = part.strip()
            project_name = re.sub(r'^[^\w&]+', '', project_name)
            project_name = re.sub(r'[^\w\s&/-]+$', '', project_name)
            
            if project_name and len(project_name) > 3:
                normalized_found = normalize_name(project_name)
                
                # Try to match with existing projects
                for project in existing_projects:
                    project_title = project.get('title', '')
                    normalized_title = normalize_name(project_title)
                    
                    # Flexible matching
                    if (normalized_found == normalized_title or 
                        normalized_found.startswith(normalized_title) or
                        normalized_title.startswith(normalized_found) or
                        normalized_found in normalized_title or
                        normalized_title in normalized_found):
                        if not any(p['title'] == project_title for p in projects):
                            projects.append({
                                'title': project_title,
                                'project': project,
                                'line': line,
                                'line_index': line_idx,
                                'found_name': project_name
                            })
                            log(f"Matched: '{project_name}' -> '{project_title}'")
                            break
    
    log(f"Total projects matched: {len(projects)}")
    return projects

def parse_project_data(text: str, project_title: str, line_index: int) -> Dict[str, str]:
    """Parse project data from text section"""
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    
    # Get section starting from the project line
    start_index = line_index
    end_index = min(len(lines), start_index + 30)
    
    # Find next "Project" line to mark end of section
    for i in range(start_index + 1, len(lines)):
        if i >= end_index:
            break
        if lines[i].strip().lower().startswith('project'):
            end_index = i
            break
    
    project_section = '\n'.join(lines[start_index:end_index])
    log(f"  Processing section ({end_index - start_index} lines)")
    
    data = {
        'bau': '',
        'area': '',
        'size': '',
        'quantity': '',
        'contractor': '',
        'architect': '',
        'client': '',
        'scope': '',
        'country': ''
    }
    
    # Patterns for each field
    patterns = {
        'bau': r'(?:bua|bau|built[-\s]*up[-\s]*area)[:\s]*([^\n]+)',
        'area': r'area[:\s]*([^\n]+)',
        'size': r'size[:\s]*([^\n]+)',
        'quantity': r'quantity[:\s]*([^\n]+)',
        'contractor': r'contractor[:\s]*([^\n]+)',
        'architect': r'architect[:\s]*([^\n]+)',
        'client': r'client[:\s]*([^\n]+)',
        'scope': r'scope[:\s]*([^\n]+)',
        'country': r'country[:\s]*([^\n]+)'
    }
    
    # Try regex patterns first
    for key, pattern in patterns.items():
        matches = re.finditer(pattern, project_section, re.IGNORECASE | re.MULTILINE)
        for match in matches:
            value = match.group(1).strip()
            # Clean up value
            value = re.sub(r'^\W+|\W+$', '', value)
            if value and len(value) > 1:
                if key == 'scope':
                    # Scope can be multi-line, continue until next label
                    start_pos = match.end()
                    scope_lines = [value]
                    remaining = project_section[start_pos:]
                    for line in remaining.split('\n')[:5]:  # Max 5 more lines
                        line = line.strip()
                        if not line:
                            continue
                        # Stop if we hit another label
                        if any(k in line.lower() for k in ['project', 'country', 'bua', 'bau', 'area', 'contractor', 'architect', 'client', 'scope', 'size', 'quantity']):
                            break
                        scope_lines.append(line)
                    value = ' '.join(scope_lines)
                data[key] = value
                break
    
    # Alternative: parse line by line
    section_lines = project_section.split('\n')
    for i, line in enumerate(section_lines):
        line_lower = line.lower()
        
        # Check for BUA (Built-Up Area) - map to bau
        if not data['bau'] and ('bua' in line_lower or ('built' in line_lower and 'area' in line_lower)):
            value = re.sub(r'(bua|built[-\s]*up[-\s]*area|built[-\s]*area)[:\s]*', '', line, flags=re.IGNORECASE).strip()
            if not value and i + 1 < len(section_lines):
                value = section_lines[i + 1].strip()
            value = re.sub(r'^\W+|\W+$', '', value)
            if value and len(value) > 1:
                data['bau'] = value
                continue
        
        for key in data.keys():
            if not data[key] and key in line_lower:
                value = re.sub(f'{key}[:\s]*', '', line, flags=re.IGNORECASE).strip()
                if not value and i + 1 < len(section_lines):
                    next_line = section_lines[i + 1].strip()
                    # Only use next line if it doesn't look like another label
                    if not any(k in next_line.lower() for k in ['project', 'country', 'bua', 'bau', 'area', 'contractor', 'architect', 'client', 'scope', 'size', 'quantity']):
                        value = next_line
                
                if key == 'scope' and value and i + 1 < len(section_lines):
                    j = i + 1
                    while j < len(section_lines):
                        next_line = section_lines[j].strip()
                        next_lower = next_line.lower()
                        if any(k in next_lower for k in ['project', 'country', 'bua', 'bau', 'area', 'contractor', 'architect', 'client', 'scope', 'size', 'quantity']):
                            break
                        if next_line:
                            value += ' ' + next_line
                        j += 1
                        if len(value) > 200:
                            break
                
                value = re.sub(r'^\W+|\W+$', '', value)
                if value and len(value) > 1:
                    data[key] = value
    
    return data

def process_pptx(pptx_path: Path, category: str, existing_projects: List[Dict]) -> Dict:
    """Process a PowerPoint file and extract project data"""
    log(f"\n=== Processing PowerPoint: {pptx_path.name} ===")
    
    # Extract text from PowerPoint
    text = extract_text_from_pptx(pptx_path)
    
    if not text or len(text.strip()) < 10:
        log(f"WARNING: No text extracted from {pptx_path.name}")
        return {'extracted': [], 'raw_text': text}
    
    # Find projects in the text
    found_projects = find_projects_in_category(text, category, existing_projects)
    
    log(f"Found {len(found_projects)} projects in PowerPoint")
    
    extracted_data = []
    for found in found_projects:
        project_title = found['title']
        line_index = found.get('line_index', 0)
        log(f"\nExtracting data for: {project_title}")
        log(f"  Found at line {line_index}: {found.get('found_name', project_title)}")
        
        extracted = parse_project_data(text, project_title, line_index)
        log(f"Extracted data: {json.dumps(extracted, indent=2, ensure_ascii=False)}")
        
        extracted_data.append({
            'project_title': project_title,
            'project': found['project'],
            'data': extracted
        })
    
    return {'extracted': extracted_data, 'raw_text': text}

def extract_project_data_from_pptx():
    """Main function to extract project data from PowerPoint files"""
    try:
        # Clear log file
        with open(LOG_FILE, 'w', encoding='utf-8') as f:
            f.write('')
        
        log('=== Starting PowerPoint-based Project Data Extraction ===')
        log(f"Base directory: {BASE_DIR}")
        
        # Load projects
        with open(PROJECTS_JSON, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        projects = data.get('projects', [])
        log(f"Loaded {len(projects)} projects")
        
        # Find PowerPoint files
        pptx_files = list(PROJECT_SLIDES_DIR.glob("*.pptx")) + list(PROJECT_SLIDES_DIR.glob("*.ppt"))
        log(f"Found {len(pptx_files)} PowerPoint files")
        
        if not pptx_files:
            log("ERROR: No PowerPoint files found in Project Slides directory")
            log("Please ensure you have .pptx or .ppt files in the Project Slides folder")
            return
        
        all_extracted = []
        total_updates = 0
        
        # Process each PowerPoint file
        for pptx_file in sorted(pptx_files):
            # Get category from filename
            category = pptx_file.stem.strip()
            log(f"\n=== Processing category: {category} ===")
            
            # Get projects in this category
            category_projects = [p for p in projects if category in p.get('categories', [])]
            log(f"Found {len(category_projects)} projects in category: {category}")
            
            if not category_projects:
                log(f"WARNING: No projects found for category '{category}'")
                continue
            
            # Process PowerPoint
            result = process_pptx(pptx_file, category, category_projects)
            all_extracted.extend(result['extracted'])
            
            # Update projects.json with extracted data
            for item in result['extracted']:
                project_data = item['data']
                project = item['project']
                
                updates = []
                for key, value in project_data.items():
                    if value and (not project.get(key) or len(str(project.get(key, ''))) < len(value)):
                        old_value = project.get(key, '')
                        project[key] = value
                        updates.append(f"{key}: {value}")
                        total_updates += 1
                
                if updates:
                    log(f"Updated {item['project_title']}: {', '.join(updates)}")
        
        # Save updated projects.json
        with open(PROJECTS_JSON, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
        log(f"\n✓ Saved updated projects.json")
        
        # Generate report
        report_lines = [
            "# PowerPoint Extraction Report",
            f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            "",
            f"PowerPoint files processed: {len(pptx_files)}",
            f"Projects updated: {len(all_extracted)}",
            f"Total field updates: {total_updates}",
            "",
            "## Extracted Data",
            ""
        ]
        
        for item in all_extracted:
            report_lines.append(f"### {item['project_title']}")
            report_lines.append("```json")
            report_lines.append(json.dumps(item['data'], indent=2, ensure_ascii=False))
            report_lines.append("```")
            report_lines.append("")
        
        with open(REPORT_FILE, 'w', encoding='utf-8') as f:
            f.write('\n'.join(report_lines))
        
        log(f"✓ Report saved to: {REPORT_FILE}")
        
        log("\n=== Extraction Complete ===")
        log(f"PowerPoint files processed: {len(pptx_files)}")
        log(f"Projects updated: {len(all_extracted)}")
        log(f"Total field updates: {total_updates}")
        log(f"\nLog saved to: {LOG_FILE}")
        log(f"Report saved to: {REPORT_FILE}")
        
    except Exception as e:
        log(f"ERROR: {e}")
        import traceback
        log(traceback.format_exc())
        raise

if __name__ == '__main__':
    extract_project_data_from_pptx()






